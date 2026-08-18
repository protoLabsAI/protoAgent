"""File previews (ADR 0092 D2): mime, clipping, office/pdf extractors, thumbnails."""

from __future__ import annotations

import logging
from pathlib import Path

from . import _config

log = logging.getLogger("protoagent.plugins.artifact")

# ── file artifacts (ADR 0092 D2) ─────────────────────────────────────────────
# save_file_artifact turns a generated file (docx/xlsx/pptx/pdf/image/text) into a
# VERSIONED artifact: bytes stored as a sidecar blob, a diffable text PREVIEW extracted
# into `code`, and — for images only — a base64 thumbnail (the cheap win; rasterizing
# Office/PDF pages needs LibreOffice/pdfium and is deferred). Every extractor imports its
# lib lazily and degrades to a note if it's missing (a lean non-desktop env without the
# doc stack), so the tool never hard-fails — it just stores the blob with a thin preview.
_PREVIEW_TRUNC = "\n… (preview truncated — download the file for the full content)"


def _guess_mime(path: Path) -> str:
    import mimetypes

    mime, _ = mimetypes.guess_type(str(path))
    return mime or "application/octet-stream"


def _clip(text: str) -> str:
    """Clip an extracted preview to _config._max_preview_bytes (utf-8), with a truncation note.
    Cuts on a codepoint boundary — back the byte cut off any trailing continuation byte
    (0b10xxxxxx) so a multi-byte char is never split (no silently-dropped straddler)."""
    text = text or ""
    data = text.encode("utf-8")
    limit = _config._max_preview_bytes()
    if len(data) <= limit:
        return text
    cut = max(0, limit - len(_PREVIEW_TRUNC.encode()))
    while cut > 0 and (data[cut] & 0xC0) == 0x80:  # inside a multi-byte sequence → back off
        cut -= 1
    return data[:cut].decode("utf-8") + _PREVIEW_TRUNC


def _extract_docx(path: Path) -> str:
    from docx import Document  # python-docx

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"# {ws.title}")
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r >= 200:  # cap rows per sheet — a preview, not the whole workbook
                out.append("… (more rows — download for all)")
                break
            out.append(", ".join("" if c is None else str(c) for c in (row or ())[:50]))
    wb.close()
    return "\n".join(out)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(str(path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"── Slide {i} ──")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return "\n".join(out)


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader  # already a core dep (pyproject)

    reader = PdfReader(str(path))
    out: list[str] = []
    for i, page in enumerate(reader.pages):
        if i >= 50:  # cap pages
            out.append("… (more pages — download for all)")
            break
        out.append(page.extract_text() or "")
    return "\n".join(out)


# ext → (extractor, degrade-note). Text-ish kinds decode verbatim in _extract_preview.
_EXTRACTORS = {
    ".docx": (_extract_docx, "python-docx"),
    ".xlsx": (_extract_xlsx, "openpyxl"),
    ".pptx": (_extract_pptx, "python-pptx"),
    ".pdf": (_extract_pdf, "pypdf"),
}
_TEXT_EXT = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".json",
    ".log",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".xml",
    ".html",
    ".htm",
    ".py",
    ".js",
    ".ts",
    ".sh",
    ".sql",
}


def _extract_preview(path: Path, data: bytes, mime: str) -> str:
    """A readable, diffable text projection of the file, capped. Office/PDF via their lib
    (lazy, degrades if absent); text/* decoded verbatim; anything else → a size note."""
    ext = path.suffix.lower()
    if ext in _EXTRACTORS:
        fn, lib = _EXTRACTORS[ext]
        try:
            return _clip(fn(path))
        except Exception:  # noqa: BLE001 — missing lib / corrupt file → thin preview, never crash
            log.debug("[artifact] %s preview extract failed", ext, exc_info=True)
            return (
                f"({ext[1:].upper()} file — no text preview: {lib} unavailable or the file "
                f"couldn't be parsed. Download it for the full content.)"
            )
    if mime.startswith("text/") or ext in _TEXT_EXT:
        return _clip(data.decode("utf-8", "replace"))
    if mime.startswith("image/"):
        return f"(image · {mime})"
    return f"(binary file · {mime} · {len(data)} bytes — download to open)"


def _thumbnail(data: bytes, mime: str) -> str | None:
    """A small base64 PNG data-URI thumbnail for IMAGE files only (Pillow, bundled on the
    desktop runtime via ADR 0092 D1). Office/PDF thumbnails need a rasterizer we don't ship
    — those show the file icon + text preview instead. Returns None on any failure."""
    if not mime.startswith("image/"):
        return None
    try:
        import base64
        from io import BytesIO

        from PIL import Image

        im = Image.open(BytesIO(data))
        im.thumbnail((320, 320))
        if im.mode not in ("RGB", "RGBA"):
            im = im.convert("RGBA")
        buf = BytesIO()
        im.save(buf, format="PNG")
        return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode()
    except Exception:  # noqa: BLE001 — Pillow absent (lean env) / undecodable → no thumb
        log.debug("[artifact] thumbnail failed", exc_info=True)
        return None
